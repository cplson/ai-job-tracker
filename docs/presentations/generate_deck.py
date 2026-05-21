#!/usr/bin/env python3
"""Generate JobTracker Capstone PowerPoint from slide definitions."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

DIR = Path(__file__).resolve().parent
OUTPUT = DIR / "JobTracker-Capstone.pptx"
# Microsoft Office "Facet" theme (designer-style geometric layout)
TEMPLATE = DIR / "designer-template.pptx"

TITLE_COLOR = RGBColor(0x1A, 0x36, 0x5D)
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
ACCENT_COLOR = RGBColor(0x2E, 0x6B, 0x9E)

LAYOUT_TITLE = 0
LAYOUT_TITLE_AND_CONTENT = 1

SLIDES = [
    {
        "title": "AI Job Tracker — Capstone",
        "subtitle": "Dunwoody Web Programming · CI/CD Focus",
        "bullets": [],
        "notes": (
            "Open with the problem: job seekers juggle many applications across spreadsheets "
            "or notes, with little feedback on whether a resume fits a role. AI Job Tracker "
            "centralizes applications, stores resumes, and uses OpenAI to analyze fit. "
            "This deck covers what we built and how we ship it safely with a Jenkins CI/CD "
            "pipeline on Linode."
        ),
        "is_title": True,
    },
    {
        "title": "What It Does",
        "bullets": [
            "Register and login with JWT authentication",
            "CRUD job applications with status workflow",
            "  Draft → Applied → Interviewing → Offer → Rejected",
            "Upload resumes (PDF, DOCX, plain text)",
            "AI analysis: summary, strengths, weaknesses, suggestions, match score",
        ],
        "notes": (
            "Walk through the user journey: sign up, add applications, upload a resume, "
            "link it to an application, run AI analysis. Demo: seed with "
            "./scripts/seed-demo-data.sh, then demo@jobtracker.app / Demo123!"
        ),
    },
    {
        "title": "Tech Stack",
        "bullets": [
            "Frontend: React 19, TypeScript, Vite 7, React Router, Bootstrap 5",
            "Backend: ASP.NET Core 8, EF Core 8, JWT, BCrypt",
            "Database: PostgreSQL 15 (Docker)",
            "AI: OpenAI Chat Completions (gpt-4.1-mini)",
            "Ops: Docker, Docker Compose, nginx, Linode Ubuntu 24.04",
        ],
        "notes": (
            "Conventional SPA plus REST API. Resume text via PdfPig and OpenXml. "
        ),
    },
    {
        "title": "Architecture",
        "bullets": [
            "Browser → nginx → React static (/var/www/jobtracker)",
            "nginx /api/ proxy → Dockerized .NET API → PostgreSQL",
            "API → OpenAI; resume files on container disk (/app/uploads)",
            "Backend layers: API → Core → Infrastructure",
        ],
        "notes": (
            "nginx serves Vite build and reverse-proxies /api/ to API on localhost:5001. "
            "API and Postgres under Docker Compose with named volume backend_pgdata."
        ),
    },
    {
        "title": "Security & API Design",
        "bullets": [
            "BCrypt password hashing; JWT Bearer on protected routes",
            "User-scoped data via /me endpoints and authorization",
            "DTOs, pagination, search, and sort on list endpoints",
            "Secrets in backend/.env — not committed to git",
        ],
        "notes": (
            "Swagger in Development only. Bridges to CI/CD: pipeline runs SAST, DAST, "
            "and container scans because we handle credentials and user data."
        ),
    },
    {
        "title": "CI/CD Philosophy",
        "bullets": [
            "Jenkins (CI/CD server): runs our build pipeline automatically on each trigger",
            "Not using GitHub Actions — the full pipeline is on a Linode server",
            "Pipeline defined in backend/Jenkinsfile — version-controlled like app code",
            "Every successful run deploys to production — quality and security checks run first",
        ],
        "notes": (
            "Full pipeline ownership for capstone: provision server, Jenkins agent, "
            "checkout through deploy on one machine. Trade-off: less PR-level CI "
            "in exchange for realistic small-team production flow."
        ),
    },
    {
        "title": "Infrastructure & Jenkins Agent",
        "bullets": [
            "Linode: cloud Linux server — hosts the live app and Jenkins build agent",
            "Jenkins agent: worker that runs pipeline steps (needs Docker installed)",
            "One-time setup scripts install Docker, .NET 8, Node 20, and nginx",
            "GitHub: stores source code; Jenkins clones main branch every build",
        ],
        "notes": (
            "Provision VM → bootstrap → register Jenkins agent (labels: jobtracker, docker, dotnet). "
            "Agent needs Docker for images, ZAP compose, and Trivy."
        ),
        "diagram": "GitHub → Jenkins → Linode agent (build, scan, and deploy on same machine)",
    },
    {
        "title": "Pipeline Overview — 8 Stages",
        "bullets": [
            "1. Checkout (Jenkins + GitHub) — download latest source code",
            "2. SonarQube — run unit tests; find code smells, bugs, and coverage gaps",
            "3. Snyk — scan source code for security vulnerabilities (SAST)",
            "4. OWASP ZAP — attack the running app for web flaws (DAST) — blocks deploy",
            "5. Trivy — scan the container image for known CVEs",
            "6. Docker Build — package the .NET API into a container image",
            "7. Build Frontend — compile React into static files for nginx",
            "8. Deploy — start containers, publish the site, reload nginx",
        ],
        "notes": (
            "Spine of the presentation. Walk through each stage left to right. "
            "Say which are third-party vs our scripts. Post-deploy: docker image prune -f."
        ),
        "font_size": 16,
    },
    {
        "title": "SonarQube — Code Quality & Tests",
        "bullets": [
            "SonarQube (self-hosted, open source): code quality and maintainability dashboard",
            "Detects code smells, bugs, duplicated code, and overly complex methods",
            "Runs our automated API tests (xUnit) — login, users, resume parsing, etc.",
            "Measures test coverage — highlights code with no tests",
            "Soft-fail today: pipeline continues even if the quality gate fails",
        ],
        "notes": (
            "JobTracker.Tests: UsersTests, ResumeTextExtractorTests. Tests run inside "
            "Sonar stage, not a separate Jenkins stage. Sonar via docker-compose.sonar.yml."
        ),
    },
    {
        "title": "Security Scanning — Three Tools",
        "bullets": [
            "Snyk (third-party SaaS): scans C# source before deploy — SQL injection patterns, hardcoded secrets, unsafe APIs",
            "OWASP ZAP (open source): crawls the API and UI like a hacker — XSS, misconfigurations, exposed endpoints",
            "Trivy (open source): scans the Docker image for known CVEs in packages and the OS — HIGH/CRITICAL only",
            "Defense in depth: source code → running application → container image",
            "Only ZAP, Docker build, frontend build, and deploy can stop a release",
        ],
        "notes": (
            "Three layers: source (Snyk), running app (ZAP), container (Trivy). "
            "ZAP uses separate compose project to avoid prod DB. ZAP may use ephemeral .env if missing."
        ),
    },
    {
        "title": "Build Artifacts",
        "bullets": [
            "Docker: turns the .NET API into a portable image (jobtracker-api:latest)",
            "Image runs as a non-root user — reduces risk if the container is compromised",
            "Vite: compiles React + TypeScript into static HTML/JS/CSS in frontend/dist/",
            "Docker Compose: starts API + PostgreSQL together with persistent database storage",
        ],
        "notes": (
            "SDK publish stage, aspnet runtime, non-root appuser. "
            "Duplicate build is a simplification to refactor later."
        ),
    },
    {
        "title": "Deploy Stage",
        "bullets": [
            "deploy.sh: our script that publishes a new version to the Linode server",
            "Starts API and PostgreSQL containers; waits until database and API are healthy",
            "Copies the React build to /var/www/jobtracker for nginx to serve",
            "nginx: reverse proxy — serves the website and forwards /api to the API container",
            "Database volume is never deleted — user data survives redeploys",
        ],
        "notes": (
            "Schema via EnsureCreated at API startup. Compose retry without volume deletion. "
            "Continuous delivery from Jenkins to live URL on same host."
        ),
    },
    {
        "title": "Secrets & Environments",
        "bullets": [
            "Jenkins stores integration passwords: GitHub access, Snyk API token, SonarQube token",
            "Server .env file stores app secrets: database password, JWT key, OpenAI key",
            "Secrets are never committed to git — only .env.example templates in the repo",
            "Single environment today: every pipeline run targets main and deploys to production",
        ],
        "notes": (
            "Jenkins holds tokens; Linode holds .env. "
            "Pipeline always targets main. Future: PR checks, staging environment."
        ),
    },
]


def _set_font(run, size_pt, bold=False, color=None, use_template: bool = False):
    if use_template:
        return
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _body_placeholder(slide):
    """Return the main content placeholder (not the title)."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    return slide.placeholders[1]


def _add_title_slide(prs: Presentation, data: dict, use_template: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    slide.shapes.title.text = data["title"]
    if data.get("subtitle"):
        try:
            _body_placeholder(slide).text = data["subtitle"]
        except (KeyError, IndexError):
            pass
    if not use_template:
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                _set_font(run, 40, bold=True, color=TITLE_COLOR)
    if slide.notes_slide.notes_text_frame:
        slide.notes_slide.notes_text_frame.text = data.get("notes", "")


def _add_content_slide(prs: Presentation, data: dict, use_template: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT])
    slide.shapes.title.text = data["title"]
    if not use_template:
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                _set_font(run, 32, bold=True, color=TITLE_COLOR)

    body = _body_placeholder(slide).text_frame
    body.clear()
    body_size = data.get("font_size", 18)
    for i, bullet in enumerate(data.get("bullets", [])):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 1 if bullet.startswith("  ") else 0
        if not use_template:
            p.font.size = Pt(body_size)
            for run in p.runs:
                _set_font(run, body_size, color=BODY_COLOR, use_template=use_template)

    if data.get("diagram"):
        p = body.add_paragraph()
        p.text = data["diagram"]
        p.level = 0
        if not use_template:
            for run in p.runs:
                _set_font(run, 14, color=ACCENT_COLOR, use_template=use_template)

    if slide.notes_slide.notes_text_frame:
        slide.notes_slide.notes_text_frame.text = data.get("notes", "")


def _load_presentation() -> tuple[Presentation, bool]:
    if TEMPLATE.exists():
        return Presentation(str(TEMPLATE)), True
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, False


def build_presentation() -> Presentation:
    prs, use_template = _load_presentation()

    for slide_data in SLIDES:
        if slide_data.get("is_title"):
            _add_title_slide(prs, slide_data, use_template)
        else:
            _add_content_slide(prs, slide_data, use_template)

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT)
    theme = "Facet designer template" if TEMPLATE.exists() else "default"
    print(f"Wrote {OUTPUT} ({len(prs.slides)} slides, theme: {theme})")


if __name__ == "__main__":
    main()
